import os
from src.file_loader import load_file


def test_load_txt():
    # 创建临时txt文件
    test_content = "This is a test file."
    with open("test.txt", "w") as f:
        f.write(test_content)

    result = load_file("test.txt")

    assert result['type'] == 'txt'
    assert result['status'] == 'success'
    assert result['content'] == test_content

    # 清理
    os.remove("test.txt")


def test_load_pdf():
    # 假设tests目录下有sample.pdf
    result = load_file("tests/sample.pdf")

    assert result['type'] == 'pdf'
    assert result['status'] == 'success'
    assert len(result['content']) > 0


def test_unsupported_file():
    result = load_file("test.docx")

    assert result['type'] == 'error'
    assert result['status'] == 'error'
    assert "Unsupported file type" in result['content']


def test_load_docx():
    # 创建临时docx文件
    from docx import Document

    doc = Document()
    doc.add_heading('Test Document', 0)
    doc.add_paragraph('This is a test paragraph.')

    table = doc.add_table(rows=2, cols=2)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = 'Name'
    hdr_cells[1].text = 'Age'
    row_cells = table.rows[1].cells
    row_cells[0].text = 'John'
    row_cells[1].text = '30'

    doc.save('test.docx')

    result = load_file('test.docx')

    assert result['type'] == 'docx'
    assert result['status'] == 'success'
    assert 'Test Document' in result['content']
    assert 'John' in result['content']
    assert '30' in result['content']

    # 验证元数据
    assert 'author' in result['metadata']
    assert 'created' in result['metadata']

    # 清理
    os.remove('test.docx')


def test_load_excel():
    # 创建临时Excel文件
    import pandas as pd

    # 创建测试数据
    data = {
        'Name': ['Alice', 'Bob', 'Charlie'],
        'Age': [25, 30, 35],
        'Salary': [5000, 6000, 7000]
    }
    df = pd.DataFrame(data)

    # 创建第二个工作表
    data2 = {
        'Product': ['Apple', 'Banana'],
        'Price': [2.5, 1.8]
    }
    df2 = pd.DataFrame(data2)

    # 写入Excel文件
    with pd.ExcelWriter('test.xlsx') as writer:
        df.to_excel(writer, sheet_name='Employees', index=False)
        df2.to_excel(writer, sheet_name='Products', index=False)

    result = load_file('test.xlsx')

    assert result['type'] == 'excel'
    assert result['status'] == 'success'
    assert 'Employees' in result['content']
    assert 'Alice' in result['content']
    assert 'Apple' in result['content']

    # 验证元数据
    assert 'sheet_names' in result['metadata']
    assert len(result['metadata']['sheet_names']) == 2
    assert 'Employees' in result['metadata']['sheet_names']

    # 清理
    os.remove('test.xlsx')


def test_load_ppt():
    # 创建临时PPT文件
    from pptx import Presentation

    prs = Presentation()

    # 标题幻灯片
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "测试演示文稿"
    subtitle.text = "这是一个测试副标题"

    # 内容幻灯片
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "第一部分"
    tf = body.text_frame
    tf.text = "这是第一点"
    p = tf.add_paragraph()
    p.text = "这是第二点"

    # 添加表格
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "表格示例"
    rows = 2
    cols = 2
    left = top = 1
    width = 8
    height = 1
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    table.cell(0, 0).text = "姓名"
    table.cell(0, 1).text = "年龄"
    table.cell(1, 0).text = "张三"
    table.cell(1, 1).text = "25"

    prs.save('test.pptx')

    result = load_file('test.pptx')

    assert result['type'] == 'ppt'
    assert result['status'] == 'success'
    assert '测试演示文稿' in result['content']
    assert '张三' in result['content']
    assert '25' in result['content']

    # 验证元数据
    assert 'title' in result['metadata']
    assert 'num_slides' in result['metadata']
    assert result['metadata']['num_slides'] == 3

    # 清理
    os.remove('test.pptx')
