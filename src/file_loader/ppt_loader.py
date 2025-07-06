import logging
from pptx import Presentation
from typing import List, Dict, Any


def load_ppt(file_path: str) -> Dict[str, Any]:
    try:
        prs = Presentation(file_path)
        result = {
            'content': '',
            'metadata': {
                'title': prs.core_properties.title or '',
                'author': prs.core_properties.author or '',
                'created': str(prs.core_properties.created) if prs.core_properties.created else '',
                'num_slides': len(prs.slides)
            },
            'slides': []
        }

        all_slides_content = []

        for i, slide in enumerate(prs.slides):
            slide_content = [f"[幻灯片 {i + 1}]"]

            # 提取标题
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                slide_title = title_shape.text
                slide_content.append(f"标题: {slide_title}")

            # 提取正文内容
            body_shapes = [shape for shape in slide.shapes if shape.has_text_frame and shape != title_shape]
            for shape in body_shapes:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        slide_content.append(run.text)

            # 提取表格
            for shape in slide.shapes:
                if shape.has_table:
                    table = shape.table
                    table_content = []
                    for row_idx, row in enumerate(table.rows):
                        row_data = []
                        for cell in row.cells:
                            row_data.append(cell.text)
                        table_content.append(f"行{row_idx + 1}: " + "\t".join(row_data))
                    slide_content.append("[表格]\n" + "\n".join(table_content))

            slide_text = "\n".join(slide_content)
            all_slides_content.append(slide_text)
            result['slides'].append({
                'index': i + 1,
                'title': slide_title if title_shape else '',
                'content': slide_text
            })

        result['content'] = "\n\n\n".join(all_slides_content)
        result['size'] = len(result['content'])
        result['status'] = 'success'

        return result

    except Exception as e:
        logging.error(f"Error loading PPT {file_path}: {str(e)}")
        return {
            'content': f"Error: {str(e)}",
            'metadata': {},
            'slides': [],
            'size': 0,
            'status': 'error'
        }
